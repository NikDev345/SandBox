from __future__ import annotations

from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation


class DocumentParser:
    """
    Universal document parser.

    Supports:
        - PDF
        - DOCX
        - PPTX
        - TXT
        - Markdown
    """

    SUPPORTED_EXTENSIONS = {
        ".pdf",
        ".docx",
        ".pptx",
        ".txt",
        ".md",
        ".markdown",
    }

    @classmethod
    def parse(cls, file_path: str) -> str:
        """
        Parse any supported document and return extracted text.
        """

        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"{file_path} not found.")

        extension = path.suffix.lower()

        if extension not in cls.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported document type: {extension}")

        parser_map = {
            ".pdf": cls._parse_pdf,
            ".docx": cls._parse_docx,
            ".pptx": cls._parse_pptx,
            ".txt": cls._parse_txt,
            ".md": cls._parse_txt,
            ".markdown": cls._parse_txt,
        }

        return parser_map[extension](path)

    # -----------------------------------------------------
    # PDF
    # -----------------------------------------------------

    @staticmethod
    def _parse_pdf(path: Path) -> str:

        text = []

        with fitz.open(path) as pdf:

            for page in pdf:

                page_text = page.get_text("text")

                if page_text:
                    text.append(page_text)

        return "\n".join(text)

    # -----------------------------------------------------
    # DOCX
    # -----------------------------------------------------

    @staticmethod
    def _parse_docx(path: Path) -> str:

        document = Document(path)

        return "\n".join(
            paragraph.text
            for paragraph in document.paragraphs
            if paragraph.text.strip()
        )

    # -----------------------------------------------------
    # PPTX
    # -----------------------------------------------------

    @staticmethod
    def _parse_pptx(path: Path) -> str:

        presentation = Presentation(path)

        text = []

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    if shape.text.strip():
                        text.append(shape.text)

        return "\n".join(text)

    # -----------------------------------------------------
    # TXT / Markdown
    # -----------------------------------------------------

    @staticmethod
    def _parse_txt(path: Path) -> str:

        return path.read_text(
            encoding="utf-8",
            errors="ignore"
        )